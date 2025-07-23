import requests
import json
import os
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import argparse
from pathlib import Path

try:
    from PIL import Image, ImageDraw
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("PIL not available. Image fitting will be limited.")

import os
import requests
from PIL import Image, ImageDraw

class SimpleImageGenerator:
    def __init__(self):
        self.api_key = os.getenv("UNSPLASH_API_KEY")
        self.headers = {'Authorization': f'Client-ID {self.api_key}'}
        self.request_count = 0  

    def generate_images(self, topic, num_slides=5):  
        """Get ALL images in ONE request - compatible with your existing code"""
        self.request_count = 0
        try:
            # SINGLE API REQUEST
            self.request_count += 1
            response = requests.get(
                "https://api.unsplash.com/search/photos",
                headers=self.headers,
                params={
                    'query': topic,
                    'per_page': num_slides,  # Now using correct parameter
                    'orientation': 'landscape'
                },
                timeout=20
            )
            response.raise_for_status()

            os.makedirs("presentation_images", exist_ok=True)
            images = []
            
            for photo in response.json()['results'][:num_slides]:
                img_id = photo['id']
                filename = f"{topic[:20]}_{img_id}.jpg"
                filepath = os.path.join("presentation_images", filename)
                
                if not os.path.exists(filepath):
                    with open(filepath, 'wb') as f:
                        f.write(requests.get(photo['urls']['regular']).content)
                
                images.append({
                    'filepath': filepath,
                    'photographer': photo['user']['name']
                })
            
            print(f"Generated {len(images)} images in 1 request")
            return images
            
        except Exception as e:
            print(f"Error: {str(e)}")
            return self._create_placeholders(topic, num_slides)

    def _create_placeholders(self, topic, num_slides):
        """Fallback image generation"""
        return [{
            'filepath': None,
            'photographer': "Placeholder"
        } for _ in range(num_slides)]


def generate_ai_content(topic, num_slides, api_key=None):
    """Generate presentation content using AI with improved error handling"""
    if api_key is None:
        api_key = os.getenv("OPENROUTER_API_KEY")
    prompt = f"""Create a PowerPoint presentation with exactly {num_slides} slides about "{topic}".

Return ONLY a valid JSON array of slide objects. Each slide must have:
- "title": A clear, descriptive slide title (5-8 words)
- "content": An array of exactly 3-4 bullet points (each 10-12 words)

Format example:
[
  {{
    "title": "Introduction to {topic}",
    "content": [
      "First comprehensive bullet point about the main concept",
      "Second detailed bullet point explaining key features",
      "Third informative bullet point covering implementation",
      "Fourth conclusive bullet point about applications"
    ]
  }}
]

Make the content educational, specific, and valuable. Focus on {topic}."""

    try:
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
            "HTTP-Referer": os.getenv("HTTP_REFERER", "https://www.webstylepress.com"),
            "X-Title": os.getenv("X_TITLE", "PowerPoint Generator")
        }

        payload = {
            "model": "anthropic/claude-3-haiku",
            "messages": [
                {
                    "role": "system", 
                    "content": "You are a professional presentation creator. Return only valid JSON array format. No additional text or explanations."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            "temperature": 0.3,
            "max_tokens": 1000
        }

        print("\nGenerating AI content...")
        response = requests.post(
            os.getenv("OPENROUTER_ENDPOINT", "https://openrouter.ai/api/v1/chat/completions"),
            headers=headers,
            json=payload,
            timeout=60
        )

        if response.status_code != 200:
            print(f"API Error: {response.status_code} - {response.text}")
            return create_default_content(topic, num_slides)

        ai_response = response.json()['choices'][0]['message']['content']
        print(f"AI Response received: {len(ai_response)} characters")
        
        # Clean response
        ai_response = ai_response.replace('```json', '').replace('```', '').strip()
        
        try:
            content = json.loads(ai_response)
            if isinstance(content, list) and len(content) > 0:
                content = content[:num_slides]
                while len(content) < num_slides:
                    content.append(create_default_slide(f"Additional Topic {len(content) + 1}", topic))
                print(f"Successfully parsed {len(content)} slides from AI")
                return content
            else:
                print("AI response is not a valid list")
                return create_default_content(topic, num_slides)
        except json.JSONDecodeError as e:
            print(f"JSON parsing error: {e}")
            print(f"Raw response: {ai_response[:200]}...")
            return parse_text_to_slides(ai_response, num_slides, topic)

    except Exception as e:
        print(f"Error generating AI content: {str(e)}")
        return create_default_content(topic, num_slides)

def parse_text_to_slides(text, num_slides, topic):
    """Parse text response into slide structure"""
    slides = []
    lines = text.split('\n')
    current_slide = None
    
    for line in lines:
        line = line.strip()
        if not line or line.startswith('[') or line.startswith(']'):
            continue
            
        if ('"title"' in line or line.startswith(('Slide', '#', '##')) or 
            (len(line) < 80 and not line.startswith(('-', '•', '*')))):
            if current_slide and len(current_slide.get('content', [])) > 0:
                slides.append(current_slide)
            clean_title = line.replace('"title":', '').replace('"', '').replace(',', '').strip()
            if clean_title:
                current_slide = {'title': clean_title, 'content': []}
        elif line.startswith(('-', '•', '*', '"')) and current_slide:
            clean_point = line.replace('"', '').replace(',', '').strip()
            for prefix in ['-', '•', '*']:
                clean_point = clean_point.replace(prefix, '', 1).strip()
            if len(clean_point) > 10 and len(current_slide['content']) < 4:
                current_slide['content'].append(clean_point)
    
    if current_slide and len(current_slide.get('content', [])) > 0:
        slides.append(current_slide)
    
    while len(slides) < num_slides:
        slides.append(create_default_slide(f"Additional Topic {len(slides) + 1}", topic))
    
    return slides[:num_slides]

def create_default_content(topic, num_slides):
    """Create default content if AI fails"""
    print("\nUsing fallback content generation...")
    
    default_slides = [
        {
            'title': f"Introduction to {topic}",
            'content': [
                f"Understanding {topic} is essential for modern applications",
                f"Key principles provide foundation for effective implementation",
                f"This presentation explores theoretical and practical aspects",
                f"Examples demonstrate value across various industries"
            ]
        },
        {
            'title': f"Key Features of {topic}",
            'content': [
                f"Core components include systematic approaches and methodologies",
                f"Advanced capabilities encompass automation and optimization",
                f"Quality assurance ensures reliability and standards compliance",
                f"Customization options provide flexibility for specific needs"
            ]
        },
        {
            'title': f"Implementation of {topic}",
            'content': [
                f"Strategic deployment minimizes risks and ensures smooth transition",
                f"Resource planning and timeline management are critical",
                f"Training programs help develop necessary skills",
                f"Continuous monitoring enables iterative improvements"
            ]
        },
        {
            'title': f"Benefits of {topic}",
            'content': [
                f"Increased efficiency through streamlined processes",
                f"Enhanced decision-making with data-driven insights",
                f"Cost reduction through optimized resource utilization",
                f"Improved scalability to adapt to changing requirements"
            ]
        },
        {
            'title': f"Best Practices for {topic}",
            'content': [
                f"Establish clear objectives and success criteria",
                f"Maintain regular stakeholder communication",
                f"Implement robust testing procedures",
                f"Document processes thoroughly for knowledge transfer"
            ]
        },
        {
            'title': f"Future of {topic}",
            'content': [
                f"Emerging technologies will shape future applications",
                f"AI integration will enhance automation capabilities",
                f"Industry standards will influence development priorities",
                f"Market demands will drive innovation"
            ]
        }
    ]
    
    return default_slides[:num_slides]

def create_default_slide(title, topic):
    """Create a single default slide"""
    return {
        'title': title,
        'content': [
            f"Comprehensive analysis reveals multiple success factors",
            f"Implementation requires careful consideration of capabilities",
            f"Best practices emphasize continuous improvement",
            f"Metrics should include quantitative and qualitative measures"
        ]
    }

def get_color_theme():
    """Get enhanced professional color themes with unique typography"""
    THEMES = {
        'modern_minimalist': {
            'primary': RGBColor(45, 55, 72),
            'secondary': RGBColor(74, 85, 104),
            'accent': RGBColor(56, 178, 172),
            'text': RGBColor(26, 32, 44),
            'light_text': RGBColor(74, 85, 104),
            'background': RGBColor(247, 250, 252),
            'content_bg': RGBColor(237, 242, 247),
            'card_bg': RGBColor(255, 255, 255),
            'border': RGBColor(203, 213, 224),
            'header_bg': RGBColor(45, 55, 72),
            'name': 'Modern Minimalist',
            'font_primary': 'Montserrat',
            'font_secondary': 'Open Sans'
        },
        'executive_elegance': {
            'primary': RGBColor(31, 41, 55),
            'secondary': RGBColor(55, 65, 81),
            'accent': RGBColor(217, 119, 6),
            'text': RGBColor(17, 24, 39),
            'light_text': RGBColor(75, 85, 99),
            'background': RGBColor(249, 250, 251),
            'content_bg': RGBColor(243, 244, 246),
            'card_bg': RGBColor(255, 255, 255),
            'border': RGBColor(209, 213, 219),
            'header_bg': RGBColor(31, 41, 55),
            'name': 'Executive Elegance',
            'font_primary': 'Playfair Display',
            'font_secondary': 'Source Sans Pro'
        },
        'tech_innovation': {
            'primary': RGBColor(79, 70, 229),
            'secondary': RGBColor(99, 102, 241),
            'accent': RGBColor(16, 185, 129),
            'text': RGBColor(31, 41, 55),
            'light_text': RGBColor(75, 85, 99),
            'background': RGBColor(240, 242, 255),
            'content_bg': RGBColor(224, 231, 255),
            'card_bg': RGBColor(255, 255, 255),
            'border': RGBColor(196, 181, 253),
            'header_bg': RGBColor(67, 56, 202),
            'name': 'Tech Innovation',
            'font_primary': 'Roboto',
            'font_secondary': 'Inter'
        }
    }
    
    theme_key = random.choice(list(THEMES.keys()))
    return THEMES[theme_key]

def fit_image_to_shape(image_path, target_width, target_height):
    """Resize and crop image to fit the target dimensions"""
    if not PIL_AVAILABLE:
        print("PIL not available, using original image")
        return image_path
        
    try:
        with Image.open(image_path) as img:
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            target_width_px = int(target_width * 96)
            target_height_px = int(target_height * 96)
            
            img_ratio = img.width / img.height
            target_ratio = target_width_px / target_height_px
            
            if img_ratio > target_ratio:
                new_height = img.height
                new_width = int(new_height * target_ratio)
                left = (img.width - new_width) // 2
                top = 0
                right = left + new_width
                bottom = img.height
            else:
                new_width = img.width
                new_height = int(new_width / target_ratio)
                left = 0
                top = (img.height - new_height) // 2
                right = img.width
                bottom = top + new_height
            
            cropped_img = img.crop((left, top, right, bottom))
            resized_img = cropped_img.resize((target_width_px, target_height_px), Image.Resampling.LANCZOS)
            
            processed_path = image_path.replace('.jpg', '_fitted.jpg')
            resized_img.save(processed_path, 'JPEG', quality=95)
            
            return processed_path
    except Exception as e:
        print(f"Error processing image {image_path}: {str(e)}")
        return image_path

def create_thank_you_slide(prs, COLORS, topic):
    """Create a professional thank you slide"""
    thank_you_slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = thank_you_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        Inches(13.33), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['background']
    background.line.fill.background()

    accent_shape = thank_you_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        Inches(4), Inches(7.5)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = COLORS['primary']
    accent_shape.line.fill.background()

    overlay = thank_you_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(3.5), Inches(0), 
        Inches(1), Inches(7.5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = COLORS['secondary']
    overlay.fill.transparency = 0.3
    overlay.line.fill.background()

    thank_you_box = thank_you_slide.shapes.add_textbox(
        Inches(4.8), Inches(2.5), Inches(8), Inches(1.5)
    )
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "THANK YOU"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.color.rgb = COLORS['text']
    thank_you_para.font.size = Pt(54)
    thank_you_para.font.bold = True
    thank_you_para.font.name = COLORS['font_primary']
    thank_you_para.alignment = PP_ALIGN.LEFT

    subtitle_box = thank_you_slide.shapes.add_textbox(
        Inches(4.8), Inches(4.2), Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"Questions & Discussion"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.color.rgb = COLORS['light_text']
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.name = COLORS['font_secondary']
    subtitle_para.alignment = PP_ALIGN.LEFT

    topic_box = thank_you_slide.shapes.add_textbox(
        Inches(4.8), Inches(5.5), Inches(8), Inches(0.8)
    )
    topic_frame = topic_box.text_frame
    topic_frame.text = f"Presentation on: {topic}"
    topic_para = topic_frame.paragraphs[0]
    topic_para.font.color.rgb = COLORS['light_text']
    topic_para.font.size = Pt(16)
    topic_para.font.name = COLORS['font_secondary']
    topic_para.alignment = PP_ALIGN.LEFT

    accent_line = thank_you_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.8), Inches(4), 
        Inches(3), Inches(0.08)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = COLORS['accent']
    accent_line.line.fill.background()

def create_powerpoint(content, topic):
    """Create professional PowerPoint presentation with improved error handling"""
    try:
        image_generator = SimpleImageGenerator()
        COLORS = get_color_theme()
        
        print(f"\nCreating presentation with {COLORS['name']} theme...")
        print(f"Typography: {COLORS['font_primary']} (Headers) + {COLORS['font_secondary']} (Body)")

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # --- EFFICIENT IMAGE FETCHING LOGIC ---
        print(f"Fetching {len(content)} images for the topic '{topic}' in a single request...")
        all_images = image_generator.generate_images(topic, num_slides=len(content))

        slide_images = {}
        for idx, slide_data in enumerate(content):
            
            if all_images and idx < len(all_images) and all_images[idx]['filepath']:
                slide_images[idx] = all_images[idx]['filepath']
                print(f"Assigned image for slide {idx + 1}: {slide_data['title']}")
            else:
                slide_images[idx] = None
                print(f"Using placeholder for slide {idx + 1}: {slide_data['title']}")
        # TITLE SLIDE
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        background = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), 
            Inches(13.33), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = COLORS['background']
        background.line.fill.background()

        accent_shape = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), 
            Inches(4), Inches(7.5)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = COLORS['primary']
        accent_shape.line.fill.background()

        overlay = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(3.5), Inches(0), 
            Inches(1), Inches(7.5)
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = COLORS['secondary']
        overlay.fill.transparency = 0.3
        overlay.line.fill.background()

        title_box = title_slide.shapes.add_textbox(
            Inches(4.8), Inches(2.5), Inches(8), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = topic.upper()
        title_para = title_frame.paragraphs[0]
        title_para.font.color.rgb = COLORS['text']
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.name = COLORS['font_primary']
        title_para.alignment = PP_ALIGN.LEFT

        subtitle_box = title_slide.shapes.add_textbox(
            Inches(4.8), Inches(4.8), Inches(8), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Professional Presentation • {datetime.now().strftime('%B %Y')}"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.color.rgb = COLORS['light_text']
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.name = COLORS['font_secondary']
        subtitle_para.alignment = PP_ALIGN.LEFT

        accent_line = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.8), Inches(4.3), 
            Inches(3), Inches(0.08)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = COLORS['accent']
        accent_line.line.fill.background()

        # CONTENT SLIDES
        for idx, slide_data in enumerate(content, 1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            background = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(0), Inches(0), 
                Inches(13.33), Inches(7.5)
            )
            background.fill.solid()
            background.fill.fore_color.rgb = COLORS['background']
            background.line.fill.background()

            content_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(0.4), Inches(0.4), 
                Inches(12.53), Inches(6.7)
            )
            content_card.fill.solid()
            content_card.fill.fore_color.rgb = COLORS['card_bg']
            content_card.line.color.rgb = COLORS['border']
            content_card.line.width = Pt(1)

            header_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(0.4), Inches(0.4), 
                Inches(12.53), Inches(1.1)
            )
            header_bg.fill.solid()
            header_bg.fill.fore_color.rgb = COLORS['header_bg']
            header_bg.line.fill.background()

            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.6), Inches(11.73), Inches(0.7)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_frame.margin_left = Inches(0.2)
            title_frame.margin_right = Inches(0.2)
            title_frame.margin_top = Inches(0.1)
            title_frame.margin_bottom = Inches(0.1)
            
            title_para = title_frame.paragraphs[0]
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.font.size = Pt(24)
            title_para.font.bold = True
            title_para.font.name = COLORS['font_primary']
            title_para.alignment = PP_ALIGN.LEFT

            content_box = slide.shapes.add_textbox(
                Inches(0.9), Inches(1.8), Inches(6.8), Inches(4.5)
            )
            content_frame = content_box.text_frame
            content_frame.clear()
            content_frame.word_wrap = True
            content_frame.margin_left = Inches(0.2)
            content_frame.margin_top = Inches(0.2)
            content_frame.margin_right = Inches(0.2)
            content_frame.margin_bottom = Inches(0.2)

            for i, point in enumerate(slide_data['content'][:4]):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                
                p.text = f"• {point}"
                p.level = 0
                p.font.size = Pt(14)
                p.font.color.rgb = COLORS['text']
                p.font.name = COLORS['font_secondary']
                p.space_after = Pt(12)
                p.line_spacing = 1.3
                p.alignment = PP_ALIGN.LEFT
                p.font.bold = False

            slide_idx = idx - 1
            if slide_idx in slide_images and slide_images[slide_idx] and os.path.exists(slide_images[slide_idx]):
                try:
                    processed_image = fit_image_to_shape(slide_images[slide_idx], 4.5, 3.8)
                    
                    pic = slide.shapes.add_picture(
                        processed_image,
                        Inches(8.2), Inches(1.8), 
                        Inches(4.5), Inches(3.8)
                    )
                    print(f"Added image to slide {idx}")
                except Exception as e:
                    print(f"Error adding image to slide {idx}: {str(e)}")
                    add_professional_image_placeholder(slide, COLORS)
            else:
                add_professional_image_placeholder(slide, COLORS)

            accent_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.9), Inches(1.6), 
                Inches(2), Inches(0.05)
            )
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = COLORS['accent']
            accent_line.line.fill.background()

            page_box = slide.shapes.add_textbox(
                Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4)
            )
            page_frame = page_box.text_frame
            page_frame.text = f"{idx}"
            page_para = page_frame.paragraphs[0]
            page_para.font.size = Pt(11)
            page_para.font.color.rgb = COLORS['light_text']
            page_para.font.name = COLORS['font_secondary']
            page_para.alignment = PP_ALIGN.CENTER

        create_thank_you_slide(prs, COLORS, topic)

        filename = f"{topic.replace(' ', '_').replace('/', '_')}_presentation_{COLORS['name'].replace(' ', '_').lower()}.pptx"
        prs.save(filename)
        
        print(f"\nProfessional PowerPoint created successfully: {filename}")
        print(f"Theme: {COLORS['name']} with {COLORS['font_primary']} + {COLORS['font_secondary']} typography")
        return filename

    except Exception as e:
        print(f"\nError creating PowerPoint: {str(e)}")
        return None

def add_professional_image_placeholder(slide, COLORS):
    """Add professional image placeholder"""
    image_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.2), Inches(1.8), 
        Inches(4.5), Inches(3.8)
    )
    image_placeholder.fill.solid()
    image_placeholder.fill.fore_color.rgb = COLORS['background']
    image_placeholder.line.color.rgb = COLORS['border']
    image_placeholder.line.width = Pt(1)

    icon_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9.9), Inches(3.2), 
        Inches(1), Inches(1)
    )
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = COLORS['primary']
    icon_circle.line.fill.background()

    img_text_box = slide.shapes.add_textbox(
        Inches(8.7), Inches(4.5), Inches(3.5), Inches(0.5)
    )
    img_frame = img_text_box.text_frame
    img_frame.text = "Image Content"
    img_para = img_frame.paragraphs[0]
    img_para.font.size = Pt(12)
    img_para.font.color.rgb = COLORS['light_text']
    img_para.font.name = COLORS['font_secondary']
    img_para.alignment = PP_ALIGN.CENTER

def main():
    """Main entry point for command line execution"""
    parser = argparse.ArgumentParser(description='Generate PowerPoint presentations')
    parser.add_argument('--topic', type=str, required=True, help='Presentation topic')
    parser.add_argument('--slides', type=int, default=5, help='Number of slides')
    
    args = parser.parse_args()
    
    OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
    
    try:
        print(f"\nGenerating presentation about '{args.topic}' with {args.slides} slides")
        content = generate_ai_content(args.topic, args.slides, OPENROUTER_API_KEY)
        
        if not content:
            print("PPT_SUCCESS:False")
            print("ERROR: Failed to generate content")
            return
        
        filename = create_powerpoint(content, args.topic)
        
        if filename and os.path.exists(filename):
            print(f"PPT_FILE:{filename}")
            print("PPT_SUCCESS:True")
        else:
            print("PPT_SUCCESS:False")
            print("ERROR: Failed to create PowerPoint file")
            
    except Exception as e:
        print("PPT_SUCCESS:False")
        print(f"ERROR: {str(e)}")

if __name__ == "__main__":
    main()
